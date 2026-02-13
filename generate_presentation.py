#!/usr/bin/env python3
"""
Generate KBLI 2025 Compliance Executive Briefing Presentation
for PT Urban Jungle Bali - Bali Zero Advisory
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────────────
DARK_GREY = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE = RGBColor(0xE8, 0x7A, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0x80, 0x80, 0x80)
DARK_ORANGE = RGBColor(0xC0, 0x60, 0x20)
RED_ORANGE = RGBColor(0xE8, 0x3A, 0x2D)
TRANSPARENT_ORANGE = RGBColor(0x3D, 0x33, 0x2A)  # dark orange-tinted bg for boxes
BOX_BG = RGBColor(0x3A, 0x3A, 0x3A)
HIGHLIGHT_BOX = RGBColor(0x4A, 0x35, 0x20)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
LOGO_PATH = "/Users/nuzantara/Desktop/images/image.png"
OUTPUT_PATH = "/Users/nuzantara/Desktop/nuzantara/compliance_presentation_urban_jungle_bali.pptx"

FONT_NAME = "Calibri"


# ── Helper Functions ───────────────────────────────────────────────────────

def set_slide_bg(slide, color=DARK_GREY):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, left=None, top=Inches(0.3), width=Inches(1.0)):
    """Add logo to slide top-right corner."""
    if left is None:
        left = SLIDE_WIDTH - width - Inches(0.4)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, width=width)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, font_color=WHITE,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(4), font_name=FONT_NAME):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_shape_box(slide, left, top, width, height, fill_color=BOX_BG,
                  border_color=None, border_width=Pt(1.5)):
    """Add a rounded rectangle shape box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color=ORANGE):
    """Add a simple rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, color=ORANGE, thickness=Pt(2)):
    """Add a horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, thickness
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title_text, y_offset=Inches(0.4)):
    """Add a standard slide title with orange underline."""
    # Title text
    add_textbox(slide, Inches(0.7), y_offset, Inches(10), Inches(0.7),
                title_text, font_size=32, font_color=ORANGE, bold=True)
    # Orange underline
    add_line(slide, Inches(0.7), y_offset + Inches(0.65), Inches(2.5), ORANGE, Pt(3))


def add_bullet_item(text_frame, text, font_size=15, font_color=WHITE,
                    bold=False, bullet_char="\u25CF", indent=0):
    """Add a bullet point item."""
    p = text_frame.add_paragraph()
    if bullet_char:
        run1 = p.add_run()
        run1.text = f"{bullet_char}  "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = ORANGE
        run1.font.name = FONT_NAME
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(font_size)
    run2.font.color.rgb = font_color
    run2.font.bold = bold
    run2.font.name = FONT_NAME
    p.space_before = Pt(6)
    p.space_after = Pt(4)
    p.alignment = PP_ALIGN.LEFT
    return p


def add_text_in_shape(shape, text, font_size=14, font_color=WHITE, bold=False,
                      alignment=PP_ALIGN.CENTER):
    """Add text inside a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = FONT_NAME
    return tf


def add_multi_text_shape(shape, lines, alignment=PP_ALIGN.CENTER):
    """Add multiple lines of text inside a shape with formatting.
    lines: list of tuples (text, font_size, font_color, bold)
    """
    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, font_size, font_color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_NAME
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return tf


def add_arrow_shape(slide, left, top, width, height, color=ORANGE):
    """Add a right arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide Builders ─────────────────────────────────────────────────────────

def build_slide_1_title(prs):
    """Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Orange accent bar at left
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ORANGE)

    # Title
    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.2),
                "KBLI 2025 COMPLIANCE", font_size=52, font_color=ORANGE, bold=True)

    # Subtitle
    add_textbox(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.6),
                "Piano di Allineamento Normativo", font_size=28, font_color=WHITE)

    # Thin divider
    add_line(slide, Inches(1.2), Inches(3.5), Inches(3), ORANGE, Pt(2))

    # Company name
    add_textbox(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.5),
                "PT Urban Jungle Bali", font_size=22, font_color=WHITE, bold=True)

    # Date
    add_textbox(slide, Inches(1.2), Inches(4.4), Inches(10), Inches(0.5),
                "Febbraio 2026", font_size=18, font_color=LIGHT_GREY)

    # Advisory
    add_textbox(slide, Inches(1.2), Inches(5.0), Inches(10), Inches(0.5),
                "Bali Zero Advisory", font_size=18, font_color=ORANGE)

    # Orange bottom bar
    add_rect(slide, Inches(0), SLIDE_HEIGHT - Inches(0.12), SLIDE_WIDTH, Inches(0.12), ORANGE)

    # Logo
    add_logo(slide, top=Inches(0.5), width=Inches(1.2))


def build_slide_2_situazione(prs):
    """Situazione Attuale."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Situazione Attuale PT Urban Jungle Bali")

    # Left column box
    left_box = add_shape_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(3.8),
                             BOX_BG, ORANGE, Pt(2))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "KBLI Attivi in OSS"
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(16)

    items_left = [
        ("68111", "Immobili di proprieta/locazione"),
        ("55193", "Villa (supporto)"),
    ]
    for code, desc in items_left:
        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = f"\u25CF  {code}"
        r1.font.size = Pt(16)
        r1.font.color.rgb = ORANGE
        r1.font.bold = True
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = f"  -  {desc}"
        r2.font.size = Pt(16)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p2.space_before = Pt(12)
        p2.space_after = Pt(8)

    # Right column box
    right_box = add_shape_box(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.8),
                              BOX_BG, LIGHT_GREY, Pt(1))
    tf2 = right_box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "KBLI in Atto Costitutivo (non in OSS)"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GREY
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(16)

    items_right = [
        ("68200", "Immobili su base contrattuale"),
        ("70209", "Consulenza gestionale"),
    ]
    for code, desc in items_right:
        p2 = tf2.add_paragraph()
        r1 = p2.add_run()
        r1.text = f"\u25CF  {code}"
        r1.font.size = Pt(16)
        r1.font.color.rgb = ORANGE
        r1.font.bold = True
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = f"  -  {desc}"
        r2.font.size = Pt(16)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p2.space_before = Pt(12)
        p2.space_after = Pt(8)

    # Bottom orange action box
    action_box = add_shape_box(slide, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.9),
                               HIGHLIGHT_BOX, ORANGE, Pt(2))
    add_text_in_shape(action_box,
                      "\u26A0  4 azioni necessarie per la compliance completa",
                      font_size=20, font_color=ORANGE, bold=True)


def build_slide_3_quadro(prs):
    """Quadro Normativo 2025-2026."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Quadro Normativo 2025-2026")

    regulations = [
        ("BPS Reg. 7/2025", "KBLI 2025 (in vigore dal 18/12/2025)"),
        ("PP 28/2025", "Licenze basate sul rischio"),
        ("BKPM Reg. 5/2025", "Nuovo regime PT PMA"),
        ("Coretax", "Sistema fiscale integrato (da gennaio 2026)"),
    ]

    y_start = Inches(1.6)
    for i, (reg, desc) in enumerate(regulations):
        y = y_start + Inches(i * 1.15)

        # Orange bullet circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.1), Inches(0.35), Inches(0.35)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        add_text_in_shape(circle, str(i + 1), font_size=14, font_color=WHITE, bold=True)

        # Regulation name (orange)
        add_textbox(slide, Inches(1.5), y, Inches(3.5), Inches(0.4),
                    reg, font_size=20, font_color=ORANGE, bold=True)

        # Arrow
        add_textbox(slide, Inches(5.0), y, Inches(0.5), Inches(0.4),
                    "\u2192", font_size=22, font_color=LIGHT_GREY, bold=True)

        # Description
        add_textbox(slide, Inches(5.5), y, Inches(6), Inches(0.5),
                    desc, font_size=18, font_color=WHITE)

        # Subtle line under each
        add_line(slide, Inches(1.5), y + Inches(0.55), Inches(10), LIGHT_GREY, Pt(0.5))

    # Deadline warning box
    deadline_box = add_shape_box(slide, Inches(3.5), Inches(6.0), Inches(6), Inches(0.8),
                                 RGBColor(0x5A, 0x25, 0x15), RED_ORANGE, Pt(2.5))
    add_text_in_shape(deadline_box,
                      "\u23F0  SCADENZA: 18 GIUGNO 2026",
                      font_size=24, font_color=RED_ORANGE, bold=True)


def build_slide_4_68111(prs):
    """KBLI 68111: Cosa Cambia."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "KBLI 68111 \u2192 Si Suddivide")

    # Source box
    src_box = add_shape_box(slide, Inches(0.8), Inches(2.2), Inches(3.2), Inches(1.5),
                            BOX_BG, ORANGE, Pt(2))
    add_multi_text_shape(src_box, [
        ("68111", 28, ORANGE, True),
        ("(KBLI 2020)", 16, LIGHT_GREY, False),
    ])

    # Arrow
    add_arrow_shape(slide, Inches(4.3), Inches(2.7), Inches(1.0), Inches(0.5), ORANGE)

    # Target boxes
    targets = [
        ("68111", "Sviluppo immobiliare residenziale", False),
        ("68112", "Affitto/gestione immobili propri", True),
        ("68129", "Immobili non residenziali", False),
    ]

    for i, (code, desc, highlight) in enumerate(targets):
        y = Inches(1.5) + Inches(i * 1.6)
        border = ORANGE if highlight else LIGHT_GREY
        bg = HIGHLIGHT_BOX if highlight else BOX_BG
        bw = Pt(3) if highlight else Pt(1)

        box = add_shape_box(slide, Inches(5.8), y, Inches(6.5), Inches(1.2),
                            bg, border, bw)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = code
        r1.font.size = Pt(22)
        r1.font.color.rgb = ORANGE
        r1.font.bold = True
        r1.font.name = FONT_NAME
        r2 = p.add_run()
        r2.text = f"  \u2192  {desc}"
        r2.font.size = Pt(16)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p.alignment = PP_ALIGN.LEFT

        if highlight:
            p2 = tf.add_paragraph()
            p2.text = "\u2B50 PROBABILE MATCH"
            p2.font.size = Pt(13)
            p2.font.color.rgb = ORANGE
            p2.font.bold = True
            p2.font.name = FONT_NAME
            p2.alignment = PP_ALIGN.LEFT
            p2.space_before = Pt(6)

    # Note at bottom
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
                "\u2713  Tutti aperti al 100% proprieta straniera",
                font_size=16, font_color=LIGHT_GREY)


def build_slide_5_68200(prs):
    """KBLI 68200: Da 1 a 4 Codici."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "KBLI 68200 \u2192 Quattro Nuovi Codici")

    codes = [
        ("68210", "Intermediazione", False),
        ("68291", "Valutazione/Perizia", False),
        ("68292", "Gestione Residenziale", True),
        ("68299", "Altre Attivita", False),
    ]

    x_positions = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.5)]

    for i, (code, desc, highlight) in enumerate(codes):
        x = x_positions[i]
        border = ORANGE if highlight else LIGHT_GREY
        bg = HIGHLIGHT_BOX if highlight else BOX_BG
        bw = Pt(3) if highlight else Pt(1)

        box = add_shape_box(slide, x, Inches(1.6), Inches(2.7), Inches(2.2),
                            bg, border, bw)
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(28)
        p.font.color.rgb = ORANGE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.space_after = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = WHITE
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

        if highlight:
            p3 = tf.add_paragraph()
            p3.text = "RACCOMANDATO"
            p3.font.size = Pt(12)
            p3.font.color.rgb = ORANGE
            p3.font.bold = True
            p3.font.name = FONT_NAME
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(12)

    # Key info section
    info_box = add_shape_box(slide, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.5),
                             BOX_BG, LIGHT_GREY, Pt(0.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Dettagli Chiave"
    p.font.size = Pt(20)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(12)

    bullets = [
        ("Rischio:", " Medio-Alto"),
        ("PMA:", " 100% proprieta straniera"),
        ("Verifica:", " 3 giorni lavorativi"),
    ]
    for label, value in bullets:
        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = f"\u25CF  {label}"
        r1.font.size = Pt(16)
        r1.font.color.rgb = ORANGE
        r1.font.bold = True
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = value
        r2.font.size = Pt(16)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p2.space_before = Pt(6)
        p2.space_after = Pt(4)


def build_slide_6_55193(prs):
    """KBLI 55193 Villa."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "KBLI 55193 \u2192 Rinumerato a 55203")

    # Visual: 55193 -> 55203
    old_box = add_shape_box(slide, Inches(1.5), Inches(2.0), Inches(3.5), Inches(2.0),
                            BOX_BG, LIGHT_GREY, Pt(2))
    add_multi_text_shape(old_box, [
        ("55193", 36, LIGHT_GREY, True),
        ("Villa", 18, WHITE, False),
        ("KBLI 2020", 14, LIGHT_GREY, False),
    ])

    # Arrow
    add_arrow_shape(slide, Inches(5.5), Inches(2.5), Inches(1.5), Inches(0.7), ORANGE)

    new_box = add_shape_box(slide, Inches(7.8), Inches(2.0), Inches(3.5), Inches(2.0),
                            HIGHLIGHT_BOX, ORANGE, Pt(3))
    add_multi_text_shape(new_box, [
        ("55203", 36, ORANGE, True),
        ("Villa", 18, WHITE, False),
        ("KBLI 2025", 14, ORANGE, False),
    ])

    # Info bullets
    info_items = [
        "Stessa attivita, stessa classificazione",
        "Rischio: Medio-Basso",
    ]
    for i, item in enumerate(info_items):
        y = Inches(4.5) + Inches(i * 0.5)
        tb = add_textbox(slide, Inches(1.5), y, Inches(10), Inches(0.4),
                         f"\u2713  {item}", font_size=18, font_color=WHITE)

    # Action required box
    action_box = add_shape_box(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(1.0),
                               HIGHLIGHT_BOX, ORANGE, Pt(2))
    tf = action_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "AZIONE RICHIESTA:  "
    r1.font.size = Pt(18)
    r1.font.color.rgb = ORANGE
    r1.font.bold = True
    r1.font.name = FONT_NAME
    r2 = p.add_run()
    r2.text = "Incorporare nell'atto costitutivo (PP 28/2025)"
    r2.font.size = Pt(18)
    r2.font.color.rgb = WHITE
    r2.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT


def build_slide_7_retroattivita(prs):
    """La Questione Retroattivita."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Migrazione: Obbligatoria, Non Facoltativa")

    points = [
        ("\u26D4", "\"Ghost Codes\" dopo il 18/06/2026",
         "I codici KBLI 2020 non saranno piu riconosciuti dal sistema"),
        ("\u26A1", "Coretax sincronizza OSS \u2194 Fisco",
         "Allineamento automatico tra sistema licenze e sistema fiscale"),
        ("\u23F3", "OSS non ancora pronto (feb 2026)",
         "Strategia: preparare ora, eseguire quando il sistema sara aggiornato"),
    ]

    for i, (icon, title, desc) in enumerate(points):
        y = Inches(1.5) + Inches(i * 1.5)

        # Icon circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.05), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        add_text_in_shape(circle, icon, font_size=18, font_color=WHITE, bold=True)

        # Title
        add_textbox(slide, Inches(1.8), y, Inches(10), Inches(0.4),
                    title, font_size=20, font_color=ORANGE, bold=True)

        # Description
        add_textbox(slide, Inches(1.8), y + Inches(0.45), Inches(10), Inches(0.5),
                    desc, font_size=16, font_color=WHITE)

    # Warning box
    warn_box = add_shape_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.0),
                             RGBColor(0x5A, 0x25, 0x15), RED_ORANGE, Pt(2.5))
    tf = warn_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "\u26A0  "
    r1.font.size = Pt(22)
    r1.font.color.rgb = RED_ORANGE
    r1.font.name = FONT_NAME
    r2 = p.add_run()
    r2.text = "NON forzare aggiornamenti prematuri in OSS"
    r2.font.size = Pt(22)
    r2.font.color.rgb = RED_ORANGE
    r2.font.bold = True
    r2.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER


def build_slide_8_atto(prs):
    """Atto Costitutivo e AGMS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Modifica Statutaria e Assemblea")

    # Left section - Passaggi
    left_box = add_shape_box(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.0),
                             BOX_BG, ORANGE, Pt(2))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Passaggi"
    p.font.size = Pt(22)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.space_after = Pt(16)

    steps = [
        ("1", "Akta Perubahan", "Atto notarile di modifica"),
        ("2", "Approvazione KEMENKUMHAM", "Tempo stimato: 30 giorni"),
        ("3", "Aggiornamento OSS", "Allineamento codici KBLI"),
    ]

    for num, title, desc in steps:
        # Number
        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = f"  {num}.  "
        r1.font.size = Pt(20)
        r1.font.color.rgb = ORANGE
        r1.font.bold = True
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = title
        r2.font.size = Pt(18)
        r2.font.color.rgb = WHITE
        r2.font.bold = True
        r2.font.name = FONT_NAME
        p2.space_before = Pt(14)

        p3 = tf.add_paragraph()
        p3.text = f"       {desc}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = LIGHT_GREY
        p3.font.name = FONT_NAME
        p3.space_after = Pt(8)

    # Right section - AGMS
    right_box = add_shape_box(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0),
                              BOX_BG, LIGHT_GREY, Pt(1))
    tf2 = right_box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "AGMS"
    p.font.size = Pt(22)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.space_after = Pt(16)

    agms_items = [
        "Combinare bilancio + modifica AoA in unica assemblea",
        "Quorum: 2/3 capitale, 2/3 voti",
        "Alternativa: Circolare dei Soci (art. 91 UU PT)",
        "Scadenza: entro giugno 2026",
    ]

    for item in agms_items:
        p2 = tf2.add_paragraph()
        r1 = p2.add_run()
        r1.text = "\u25CF  "
        r1.font.size = Pt(16)
        r1.font.color.rgb = ORANGE
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = item
        r2.font.size = Pt(15)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p2.space_before = Pt(10)
        p2.space_after = Pt(6)


def build_slide_9_lkpm(prs):
    """LKPM: Verifica e Regolarizzazione."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "LKPM Q3 2025 - Verifica e Regolarizzazione")

    # Info shape (circle, not octagon - calmer)
    info_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.9), Inches(1.6), Inches(1.0), Inches(1.0)
    )
    info_shape.fill.solid()
    info_shape.fill.fore_color.rgb = ORANGE
    info_shape.line.fill.background()
    add_text_in_shape(info_shape, "i", font_size=32, font_color=WHITE, bold=True)

    # Scadenza info
    add_textbox(slide, Inches(2.3), Inches(1.6), Inches(9), Inches(0.5),
                "Scadenza originaria Q3: 15 ottobre 2025", font_size=20, font_color=WHITE, bold=True)

    add_textbox(slide, Inches(2.3), Inches(2.1), Inches(9), Inches(0.5),
                "Il report risulta incompleto o da revisionare", font_size=16, font_color=LIGHT_GREY)

    # Left section - Sistema graduato
    add_textbox(slide, Inches(0.9), Inches(3.0), Inches(5), Inches(0.4),
                "Sistema Sanzionatorio Graduato:", font_size=18, font_color=ORANGE, bold=True)

    graduated = [
        ("1 trimestre mancato", "Registrato dal sistema, nessuna sanzione automatica"),
        ("2 trimestri consecutivi", "Avvisi scritti (Surat Peringatan)"),
        ("Dopo ripetuti avvisi", "Sospensione temporanea (30 gg per correggere)"),
    ]
    for i, (level, desc) in enumerate(graduated):
        y = Inches(3.5) + Inches(i * 0.7)
        add_textbox(slide, Inches(1.2), y, Inches(5), Inches(0.3),
                    f"\u25B8  {level}", font_size=15, font_color=ORANGE, bold=True)
        add_textbox(slide, Inches(1.5), y + Inches(0.25), Inches(5), Inches(0.3),
                    desc, font_size=13, font_color=LIGHT_GREY)

    # Right section - Cosa sapere
    info_box = add_shape_box(slide, Inches(6.5), Inches(3.0), Inches(5.8), Inches(2.8),
                             BOX_BG, ORANGE, Pt(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Da Sapere"
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.space_after = Pt(12)

    info_items = [
        "Il portale OSS chiude la finestra dopo la scadenza",
        "I dati mancanti si accumulano nel trimestre successivo",
        "PT PMA monitorate con piu attenzione di PT locali",
        "La compliance LKPM incide sul profilo di rischio futuro",
    ]

    for item in info_items:
        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = "\u25CF  "
        r1.font.size = Pt(14)
        r1.font.color.rgb = ORANGE
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = item
        r2.font.size = Pt(14)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p2.space_before = Pt(8)
        p2.space_after = Pt(4)

    # Action box at bottom
    action_box = add_shape_box(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.8),
                               HIGHLIGHT_BOX, ORANGE, Pt(2))
    tf2 = action_box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    r1 = p.add_run()
    r1.text = "RACCOMANDAZIONE:  "
    r1.font.size = Pt(16)
    r1.font.color.rgb = ORANGE
    r1.font.bold = True
    r1.font.name = FONT_NAME
    r2 = p.add_run()
    r2.text = "Verificare stato Q3 nel sistema, regolarizzare, assicurare Q4 e Q1 in ordine"
    r2.font.size = Pt(16)
    r2.font.color.rgb = WHITE
    r2.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER


def build_slide_10_ota_compliance(prs):
    """Compliance OTA - Deadline 31 Marzo 2026."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Compliance OTA - Deadline 31 Marzo 2026")

    # Regulation info
    add_textbox(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4),
                "Circolare Kemenparekraf B/SD/80/II.01/D.3.3/2025 (8 dicembre 2025)",
                font_size=16, font_color=LIGHT_GREY)

    # Main requirement box
    req_box = add_shape_box(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.3),
                            BOX_BG, ORANGE, Pt(2))
    tf = req_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Requisiti per listing su OTA (Airbnb, Booking.com, Agoda):"
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_NAME
    p.space_after = Pt(10)

    requirements = [
        "NIB (Nomor Induk Berusaha) valido",
        "KBLI appropriato (es. 55203 Villa, 68112 gestione immobiliare)",
        "Licenze turistiche complete e aggiornate",
        "Label \"Terdaftar dan Berizin\" visibile",
    ]

    for req in requirements:
        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = "\u2713  "
        r1.font.size = Pt(16)
        r1.font.color.rgb = ORANGE
        r1.font.name = FONT_NAME
        r2 = p2.add_run()
        r2.text = req
        r2.font.size = Pt(15)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_NAME
        p2.space_before = Pt(6)
        p2.space_after = Pt(4)

    # Deadline warning box
    deadline_box = add_shape_box(slide, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.9),
                                 RGBColor(0x5A, 0x25, 0x15), RED_ORANGE, Pt(3))
    tf2 = deadline_box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    r1 = p.add_run()
    r1.text = "\u26A0  DEADLINE CRITICA: 31 MARZO 2026  "
    r1.font.size = Pt(20)
    r1.font.color.rgb = RED_ORANGE
    r1.font.bold = True
    r1.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)

    p2 = tf2.add_paragraph()
    p2.text = "Proprieta non conformi = delisting automatico dalle piattaforme"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    # Enforcement info
    add_textbox(slide, Inches(0.9), Inches(4.6), Inches(11), Inches(0.4),
                "Enforcement Governativo Bali (Governatore Koster):",
                font_size=18, font_color=ORANGE, bold=True)

    enforcement = [
        ("2,000+", "Proprieta a rischio delisting"),
        ("Focus", "Zoning, PBG/SLF, registrazione fiscale"),
        ("Crackdown", "Coordinato Provincia-Kemenparekraf-OTA"),
    ]

    for i, (label, desc) in enumerate(enforcement):
        x = Inches(0.9) + Inches(i * 4.0)
        y = Inches(5.1)

        box = add_shape_box(slide, x, y, Inches(3.7), Inches(1.0),
                           BOX_BG, LIGHT_GREY, Pt(1))
        tf3 = box.text_frame
        tf3.word_wrap = True
        p = tf3.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = ORANGE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        p2 = tf3.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.CENTER

    # Impact box
    impact_box = add_shape_box(slide, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6),
                               HIGHLIGHT_BOX, ORANGE, Pt(1.5))
    tf4 = impact_box.text_frame
    tf4.word_wrap = True
    p = tf4.paragraphs[0]
    r1 = p.add_run()
    r1.text = "IMPATTO PT URBAN JUNGLE:  "
    r1.font.size = Pt(15)
    r1.font.color.rgb = ORANGE
    r1.font.bold = True
    r1.font.name = FONT_NAME
    r2 = p.add_run()
    r2.text = "Verificare tutte le proprieta nel portfolio abbiano NIB + KBLI + licenze prima del 31/03"
    r2.font.size = Pt(15)
    r2.font.color.rgb = WHITE
    r2.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT


def build_slide_11_roadmap(prs):
    """Piano Operativo / Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Roadmap di Compliance")

    # Timeline items
    timeline = [
        ("FEB 2026", [
            ("KBLI Audit & Mapping", ORANGE, "CRITICA"),
            ("Verifica compliance OTA", RED_ORANGE, "CRITICA"),
        ]),
        ("MAR 2026", [
            ("OTA Compliance (31/03)", RED_ORANGE, "DEADLINE"),
            ("Preparazione AGMS", ORANGE, None),
            ("Modifica Atto Costitutivo", ORANGE, None),
        ]),
        ("APR 2026", [
            ("LKPM Q1 (15/04)", ORANGE, None),
            ("SPT Badan (30/04)", ORANGE, None),
        ]),
        ("MAG-GIU", [
            ("AGMS + Aggiornamento OSS", ORANGE, None),
        ]),
        ("18 GIU", [
            ("DEADLINE KBLI 2025", RED_ORANGE, "DEADLINE"),
        ]),
    ]

    y_start = Inches(1.5)
    bar_height = Inches(0.35)

    for i, (month, items) in enumerate(timeline):
        y = y_start + Inches(i * 1.1)

        # Month label
        month_box = add_shape_box(slide, Inches(0.5), y, Inches(2.0), Inches(0.9),
                                  BOX_BG, ORANGE if i < 4 else RED_ORANGE, Pt(1.5))
        add_text_in_shape(month_box, month, font_size=14, font_color=ORANGE if i < 4 else RED_ORANGE,
                          bold=True)

        # Timeline bars
        for j, (item_text, color, badge) in enumerate(items):
            bar_y = y + Inches(j * 0.42) + Inches(0.05)
            bar_width = Inches(7.5) if badge == "DEADLINE" else Inches(5.5 + j * 0.8)

            bar = add_rect(slide, Inches(2.8), bar_y, bar_width, bar_height, color)
            bar_tf = bar.text_frame
            bar_tf.word_wrap = True
            p = bar_tf.paragraphs[0]
            p.text = f"  {item_text}"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.LEFT

            if badge:
                badge_x = Inches(2.8) + bar_width + Inches(0.2)
                badge_box = add_shape_box(slide, badge_x, bar_y, Inches(1.5), bar_height,
                                          RGBColor(0x5A, 0x25, 0x15) if badge != "DEADLINE" else RED_ORANGE,
                                          color, Pt(1))
                add_text_in_shape(badge_box, badge, font_size=10, font_color=color if badge != "DEADLINE" else WHITE,
                                  bold=True)

        # Vertical timeline line
        if i < len(timeline) - 1:
            add_line(slide, Inches(1.5), y + Inches(0.9), Pt(2), ORANGE, Inches(0.2))


def build_slide_12_scadenze(prs):
    """Scadenze Chiave / Calendar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Calendario Scadenze")

    # Table data
    headers = ["Scadenza", "Data", "Status"]
    rows = [
        ("Compliance OTA", "31 marzo 2026", "DEADLINE"),
        ("LKPM Q1 2026", "15 aprile 2026", "Da preparare"),
        ("SPT Tahunan Badan", "30 aprile 2026", "Da preparare"),
        ("AGMS 2025", "Entro 30 giugno 2026", "Da programmare"),
        ("KBLI 2025 Compliance", "18 giugno 2026", "DEADLINE"),
        ("KEMENKUMHAM filing", "30 gg dopo AGMS", "Da eseguire"),
    ]

    col_widths = [Inches(4.5), Inches(3.5), Inches(3.5)]
    table_left = Inches(0.9)
    table_top = Inches(1.5)
    row_height = Inches(0.75)

    # Header row
    for j, header in enumerate(headers):
        x = table_left + sum(col_widths[:j], Inches(0))
        header_box = add_rect(slide, x, table_top, col_widths[j] - Inches(0.05),
                              Inches(0.6), ORANGE)
        add_text_in_shape(header_box, header, font_size=16, font_color=WHITE, bold=True)

    # Data rows
    for i, (scadenza, data, status) in enumerate(rows):
        y = table_top + Inches(0.65) + Inches(i * 0.85)
        bg_color = BOX_BG if i % 2 == 0 else RGBColor(0x35, 0x35, 0x35)

        # Scadenza
        x0 = table_left
        cell0 = add_rect(slide, x0, y, col_widths[0] - Inches(0.05), Inches(0.7), bg_color)
        add_text_in_shape(cell0, scadenza, font_size=15, font_color=WHITE, bold=True,
                          alignment=PP_ALIGN.LEFT)

        # Data
        x1 = table_left + col_widths[0]
        cell1 = add_rect(slide, x1, y, col_widths[1] - Inches(0.05), Inches(0.7), bg_color)
        add_text_in_shape(cell1, data, font_size=15, font_color=WHITE,
                          alignment=PP_ALIGN.CENTER)

        # Status
        x2 = table_left + col_widths[0] + col_widths[1]
        is_deadline = (status == "DEADLINE")
        status_bg = RED_ORANGE if is_deadline else bg_color
        cell2 = add_rect(slide, x2, y, col_widths[2] - Inches(0.05), Inches(0.7), status_bg)
        status_color = WHITE if is_deadline else ORANGE
        add_text_in_shape(cell2, status, font_size=15,
                          font_color=status_color,
                          bold=is_deadline, alignment=PP_ALIGN.CENTER)

    # Bottom orange bar
    add_rect(slide, Inches(0), SLIDE_HEIGHT - Inches(0.06), SLIDE_WIDTH, Inches(0.06), ORANGE)


def build_slide_13_prossimi_passi(prs):
    """Prossimi Passi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, ORANGE)

    add_title_bar(slide, "Prossimi Passi Immediati")

    steps = [
        ("1", "Verifica Compliance OTA",
         "NIB + KBLI + licenze per tutte le proprieta su Airbnb/Booking (deadline 31/03/2026)"),
        ("2", "KBLI Audit & Mapping",
         "Analisi completa codici attuali vs KBLI 2025"),
        ("3", "Preventivo Dettagliato",
         "Costi per pacchetto compliance completo"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.5) + Inches(i * 1.6)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.15), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        add_text_in_shape(circle, num, font_size=28, font_color=WHITE, bold=True)

        # Content box
        content_box = add_shape_box(slide, Inches(2.0), y, Inches(10.3), Inches(1.1),
                                    BOX_BG, ORANGE, Pt(1.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.color.rgb = ORANGE
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = WHITE
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.LEFT

    # Contact info at bottom
    contact_box = add_shape_box(slide, Inches(2.5), Inches(6.2), Inches(8), Inches(0.7),
                                HIGHLIGHT_BOX, ORANGE, Pt(1))
    add_text_in_shape(contact_box,
                      "Bali Zero Advisory  \u2014  Compliance & Regulatory",
                      font_size=18, font_color=ORANGE, bold=True)


def build_slide_14_closing(prs):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Orange accent bar at left
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, ORANGE)

    # Centered logo (larger)
    if os.path.exists(LOGO_PATH):
        logo_width = Inches(2.0)
        logo_left = (SLIDE_WIDTH - logo_width) / 2
        slide.shapes.add_picture(LOGO_PATH, int(logo_left), Inches(1.5), width=logo_width)

    # Grazie
    add_textbox(slide, Inches(0), Inches(3.5), SLIDE_WIDTH, Inches(0.9),
                "Grazie", font_size=48, font_color=ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Company name
    add_textbox(slide, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.6),
                "Bali Zero Advisory", font_size=24, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, Inches(0), Inches(5.2), SLIDE_WIDTH, Inches(0.5),
                "Compliance. Strategia. Risultati.", font_size=18, font_color=LIGHT_GREY,
                alignment=PP_ALIGN.CENTER)

    # Orange bottom bar
    add_rect(slide, Inches(0), SLIDE_HEIGHT - Inches(0.12), SLIDE_WIDTH, Inches(0.12), ORANGE)


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    # Build all slides
    print("Building Slide 1: Title...")
    build_slide_1_title(prs)

    print("Building Slide 2: Situazione Attuale...")
    build_slide_2_situazione(prs)

    print("Building Slide 3: Quadro Normativo...")
    build_slide_3_quadro(prs)

    print("Building Slide 4: KBLI 68111...")
    build_slide_4_68111(prs)

    print("Building Slide 5: KBLI 68200...")
    build_slide_5_68200(prs)

    print("Building Slide 6: KBLI 55193...")
    build_slide_6_55193(prs)

    print("Building Slide 7: Retroattivita...")
    build_slide_7_retroattivita(prs)

    print("Building Slide 8: Atto Costitutivo...")
    build_slide_8_atto(prs)

    print("Building Slide 9: LKPM...")
    build_slide_9_lkpm(prs)

    print("Building Slide 10: OTA Compliance...")
    build_slide_10_ota_compliance(prs)

    print("Building Slide 11: Roadmap...")
    build_slide_11_roadmap(prs)

    print("Building Slide 12: Scadenze...")
    build_slide_12_scadenze(prs)

    print("Building Slide 13: Prossimi Passi...")
    build_slide_13_prossimi_passi(prs)

    print("Building Slide 14: Closing...")
    build_slide_14_closing(prs)

    # Save
    prs.save(OUTPUT_PATH)
    print(f"\nPresentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
